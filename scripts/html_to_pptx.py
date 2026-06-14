"""
team_presentation.html → PowerPoint (high-fidelity layout & style)
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "team_presentation.pptx"

# ── CSS theme ──
C = {
    "primary": RGBColor(0x25, 0x63, 0xEB),
    "primary_light": RGBColor(0xDB, 0xEA, 0xFE),
    "primary_dark": RGBColor(0x1E, 0x40, 0xAF),
    "accent": RGBColor(0x05, 0x96, 0x69),
    "accent_light": RGBColor(0xD1, 0xFA, 0xE5),
    "warning": RGBColor(0xD9, 0x77, 0x06),
    "warning_light": RGBColor(0xFE, 0xF3, 0xC7),
    "danger": RGBColor(0xDC, 0x26, 0x26),
    "danger_light": RGBColor(0xFE, 0xE2, 0xE2),
    "bg": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_soft": RGBColor(0xF8, 0xFA, 0xFC),
    "bg_muted": RGBColor(0xF1, 0xF5, 0xF9),
    "border": RGBColor(0xE2, 0xE8, 0xF0),
    "text": RGBColor(0x1E, 0x29, 0x3B),
    "text_sec": RGBColor(0x47, 0x55, 0x69),
    "text_muted": RGBColor(0x94, 0xA3, 0xB8),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "green_bg": RGBColor(0xF0, 0xFD, 0xF4),
    "green_border": RGBColor(0xA7, 0xF3, 0xD0),
    "red_bg": RGBColor(0xFE, 0xF2, 0xF2),
    "red_border": RGBColor(0xFE, 0xCA, 0xCA),
    "insight_bg": RGBColor(0xEF, 0xF6, 0xFF),
    "insight_success": RGBColor(0xF0, 0xFD, 0xF4),
    "code_bg": RGBColor(0xF8, 0xFA, 0xFC),
}

FONT = "맑은 고딕"
MONO = "Consolas"
ML, MT, MR = Inches(0.55), Inches(0.42), Inches(0.55)
CW = Inches(13.333) - ML - MR


class Builder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_no = 0
        self.total = 17

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _font(self, run, size=14, bold=False, color=None, name=FONT):
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    def _rect(self, slide, l, t, w, h, fill, line=None, radius=0.04):
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line:
            sh.line.color.rgb = line
            sh.line.width = Pt(1)
        else:
            sh.line.fill.background()
        if sh.adjustments:
            sh.adjustments[0] = radius
        return sh

    def _text_in(self, shape, text, size=13, bold=False, color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(6)
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        self._font(r, size, bold, color or C["text"])
        return tf

    def _textbox(self, slide, l, t, w, h, text, size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        self._font(r, size, bold, color or C["text"])
        return box

    def _progress(self, slide):
        self.slide_no += 1
        w = int(CW * self.slide_no / self.total)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, Pt(4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C["primary"]
        bar.line.fill.background()

    def _header(self, slide, title, subtitle):
        self._textbox(slide, ML, MT, CW, Inches(0.55), title, 30, True, C["text"])
        self._textbox(slide, ML, Inches(0.95), CW, Inches(0.38), subtitle, 14, False, C["text_sec"])
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Inches(1.38), CW, Pt(2))
        ln.fill.solid()
        ln.fill.fore_color.rgb = C["bg_muted"]
        ln.line.fill.background()
        return Inches(1.55)

    def _h3(self, slide, l, t, w, text):
        return self._textbox(slide, l, t, w, Inches(0.32), text, 16, True, C["primary_dark"])

    def _bullets(self, slide, l, t, w, h, items, size=13):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(5)
            r = p.add_run()
            r.text = f"• {item}"
            self._font(r, size, False, C["text_sec"])
        return box

    def _card(self, slide, l, t, w, h, title, desc, icon=None, accent_left=None):
        self._rect(slide, l, t, w, h, C["bg_soft"], C["border"])
        y = t + Inches(0.12)
        if icon:
            self._textbox(slide, l + Inches(0.12), y, w - Inches(0.2), Inches(0.35), icon, 22, align=PP_ALIGN.LEFT)
            y += Inches(0.35)
        if accent_left:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(5), h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_left
            bar.line.fill.background()
        self._textbox(slide, l + Inches(0.14), y, w - Inches(0.22), Inches(0.28), title, 13, True)
        self._textbox(slide, l + Inches(0.14), y + Inches(0.28), w - Inches(0.22), h - Inches(0.42), desc, 11, False, C["text_sec"])

    def _insight(self, slide, l, t, w, h, text, kind="info"):
        fills = {"info": C["insight_bg"], "danger": C["red_bg"], "success": C["insight_success"], "warning": C["warning_light"]}
        borders = {"info": C["primary"], "danger": C["danger"], "success": C["accent"], "warning": C["warning"]}
        self._rect(slide, l, t, w, h, fills[kind], C["border"])
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Pt(5), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = borders[kind]
        bar.line.fill.background()
        self._textbox(slide, l + Inches(0.18), t + Inches(0.08), w - Inches(0.25), h - Inches(0.12), text, 12, False, C["text"])

    def _result_box(self, slide, l, t, w, h, label, value, value_color, bg=None, border=None):
        self._rect(slide, l, t, w, h, bg or C["bg_soft"], border or C["border"])
        self._textbox(slide, l, t + Inches(0.12), w, Inches(0.22), label.upper(), 10, True, C["text_muted"], PP_ALIGN.CENTER)
        self._textbox(slide, l, t + Inches(0.38), w, Inches(0.55), value, 30, True, value_color, PP_ALIGN.CENTER)

    def _table(self, slide, l, t, w, data, col_widths=None, row_h=Inches(0.36)):
        rows, cols = len(data), len(data[0])
        ts = slide.shapes.add_table(rows, cols, l, t, w, row_h * rows)
        tb = ts.table
        if col_widths:
            for i, cw in enumerate(col_widths):
                tb.columns[i].width = cw
        for r in range(rows):
            for c in range(cols):
                cell = tb.cell(r, c)
                cell.text = str(data[r][c])
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C["primary"]
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C["bg_soft"]
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C["bg"]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        if r == 0:
                            self._font(run, 11, True, C["white"])
                        else:
                            self._font(run, 11, False, C["text_sec"])
        return ts

    def _code_block(self, slide, l, t, w, h, lines):
        sh = self._rect(slide, l, t, w, h, C["code_bg"], C["border"])
        tf = sh.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_top = Pt(10)
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            self._font(r, 10, False, C["text"], MONO)
        return sh

    def _image_card(self, slide, l, t, w, h, img_path, caption):
        self._rect(slide, l, t, w, h, C["bg_soft"], C["border"])
        p = ROOT / img_path
        if p.exists():
            slide.shapes.add_picture(str(p), l + Inches(0.12), t + Inches(0.12), width=w - Inches(0.24), height=h - Inches(0.55))
        else:
            self._textbox(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.4), Inches(0.3), f"[{img_path}]", 10, False, C["text_muted"], PP_ALIGN.CENTER)
        self._textbox(slide, l + Inches(0.1), t + h - Inches(0.42), w - Inches(0.2), Inches(0.35), caption, 10, True, C["text_muted"], PP_ALIGN.CENTER)

    def _cover(self, title, subtitle, meta, variant="start"):
        slide = self._blank()
        self._progress(slide)
        # gradient approximation
        if variant == "start":
            self._rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(0xF0, 0xF9, 0xFF))
            self._rect(slide, Inches(8), Inches(4), Inches(5.5), Inches(3.5), RGBColor(0xF0, 0xFD, 0xF4))
        else:
            self._rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(0xF0, 0xFD, 0xF4))
            self._rect(slide, Inches(8), Inches(0), Inches(5.5), Inches(4), RGBColor(0xEF, 0xF6, 0xFF))

        self._textbox(slide, ML, Inches(2.3), CW, Inches(1.2), title, 38, True, C["primary"], PP_ALIGN.CENTER)
        self._textbox(slide, ML, Inches(3.55), CW, Inches(0.55), subtitle, 17, False, C["text_sec"], PP_ALIGN.CENTER)
        self._textbox(slide, ML, Inches(4.8), CW, Inches(0.8), meta, 13, False, C["text_muted"], PP_ALIGN.CENTER)

    def build(self):
        self._slide_01()
        self._slide_02()
        self._slide_03()
        self._slide_04()
        self._slide_05()
        self._slide_06()
        self._slide_07()
        self._slide_08()
        self._slide_09()
        self._slide_10()
        self._slide_11()
        self._slide_12()
        self._slide_13()
        self._slide_14()
        self._slide_15()
        self._slide_16()
        self._slide_17()
        self.prs.save(OUTPUT)
        return OUTPUT

    def _slide_01(self):
        self._cover(
            "딥러닝 기반 농작물 잎 이미지\n병해충 진단 시스템",
            "도메인 시프트(Domain Shift) 발견 및 극복 연구 보고",
            "딥러닝분석 2조",
        )

    def _slide_02(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "목차", "발표 구성 및 흐름 안내")
        cards = [
            ("1. 문제 정의", "농업 병해충 진단의 필요성과 프로젝트 목표"),
            ("2. 데이터셋 분석", "Kaggle 데이터셋 구성 및 전처리·증강 전략"),
            ("3. 모델 아키텍처 및 최적화", "CNN → GAP 경량화 → EfficientNet-B0 도입"),
            ("4. 실험실 모델 성과", "EfficientNet-B0 검증 정확도 99.73% 달성"),
            ("5. 도메인 시프트 발견", "야외 환경 테스트 시 정확도 19.12%로 폭락"),
            ("6. 가설 검증 및 결론", "3가지 가설 실험을 통한 최적 해결책 도출"),
        ]
        gw, gh, gap = Inches(3.95), Inches(1.35), Inches(0.18)
        for i, (t, d) in enumerate(cards):
            col, row = i % 2, i // 2
            self._card(slide, ML + col * (gw + gap), y + row * (gh + gap), gw, gh, t, d)

    def _slide_03(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "1. 문제 정의", "농업 분야에서 딥러닝 기반 병해충 자동 진단의 필요성")
        self._h3(slide, ML, y, CW, "현재의 문제 상황")
        self._bullets(slide, ML, y + Inches(0.32), CW, Inches(1.1), [
            "전 세계 농작물의 약 20~40%가 매년 병해충으로 인해 손실",
            "전문가 육안 진단에 의존 → 지리적·시간적 접근성 한계 존재",
            "비전문가(농민)가 초기 증상을 정확히 구분하기 어려워 대응 시기 지연",
        ])
        self._h3(slide, ML, y + Inches(1.5), CW, "우리의 해결 목표")
        self._card(slide, ML, y + Inches(1.85), Inches(5.95), Inches(1.15), "접근성 확보",
                   "스마트폰 촬영 이미지 1장으로 38종 병해충을 자동 분류하는 딥러닝 모델 구축")
        self._card(slide, ML + Inches(6.15), y + Inches(1.85), Inches(5.95), Inches(1.15), "실시간 진단",
                   "GPU 가속 추론 환경을 활용하여 빠른 응답 시간 내에 높은 정확도의 진단 결과 제공")
        self._insight(slide, ML, y + Inches(3.15), CW, Inches(0.85),
                      "핵심 가치: 전문 인력 없이도 모바일 디바이스를 통해 누구나 즉시 병해충 종류를 식별하고 조기 대응할 수 있는 기술적 기반을 구축한다.")

    def _slide_04(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "2-1. 데이터셋 개요", "Kaggle New Plant Diseases Dataset (Augmented)")
        bw = (CW - Inches(0.36)) / 3
        self._result_box(slide, ML, y, bw, Inches(1.05), "학습 이미지", "70,295", C["primary"])
        self._result_box(slide, ML + bw + Inches(0.18), y, bw, Inches(1.05), "검증 이미지", "17,572", C["accent"])
        self._result_box(slide, ML + 2 * (bw + Inches(0.18)), y, bw, Inches(1.05), "분류 클래스", "38", C["warning"])
        self._h3(slide, ML, y + Inches(1.2), CW, "클래스 구성 예시")
        self._table(slide, ML, y + Inches(1.52), CW, [
            ["작물", "포함 클래스", "클래스 수"],
            ["Apple", "Apple Scab, Black Rot, Cedar Rust, Healthy", "4"],
            ["Tomato", "Bacterial Spot, Early/Late Blight, Mosaic Virus, Healthy 등", "10"],
            ["Potato", "Early Blight, Late Blight, Healthy", "3"],
            ["Grape", "Black Rot, Esca, Leaf Blight, Healthy", "4"],
            ["Corn", "Cercospora, Common Rust, Northern Leaf Blight, Healthy", "4"],
        ], [Inches(1.2), Inches(8.5), Inches(1.5)])
        self._textbox(slide, ML, y + Inches(3.55), CW, Inches(0.25), "* 총 10종 작물, 38개 클래스 (정상 + 질병 포함)", 10, False, C["text_muted"])

    def _slide_05(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "2-2. 데이터 전처리 및 증강", "Albumentations 라이브러리 기반의 이미지 변환 파이프라인")
        col_w = (CW - Inches(0.3)) / 2
        lx, rx = ML, ML + col_w + Inches(0.3)
        self._h3(slide, lx, y, col_w, "공통 전처리")
        for i, (t, d) in enumerate([
            ("Resize (224×224)", "CNN 입력 해상도 통일"),
            ("Normalize (ImageNet μ, σ)", "경사 하강법 수렴 속도 향상"),
            ("ToTensorV2", "PyTorch FloatTensor로 변환"),
        ]):
            self._card(slide, lx, y + Inches(0.35) + i * Inches(0.95), col_w, Inches(0.82), t, d)
        self._h3(slide, rx, y, col_w, "학습 전용 증강  [Train Only]")
        for i, (t, d) in enumerate([
            ("RandomRotation (±30°)", "촬영 각도 변화에 대한 기하학적 불변성 학습"),
            ("HorizontalFlip (p=0.5)", "좌우 비대칭 패턴을 균등하게 학습"),
            ("ColorJitter", "밝기·대비·채도·색조 변화로 조도 강건성 확보"),
        ]):
            self._card(slide, rx, y + Inches(0.35) + i * Inches(0.95), col_w, Inches(0.82), t, d)

    def _slide_06(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "3-1. 베이스라인 CNN 아키텍처", "PlantDiseaseCNN — 4-Block Convolutional Neural Network")
        self._h3(slide, ML, y, CW, "텐서 차원 흐름 (Feature Mapping)")
        blocks = [
            ("Input", "(3,224,224)"), ("Block 1", "(32,112,112)"), ("Block 2", "(64,56,56)"),
            ("Block 3", "(128,28,28)"), ("Block 4", "(256,14,14)"), ("Classifier", "50176→512→38"),
        ]
        bw, x = Inches(1.55), ML
        for i, (bt, bd) in enumerate(blocks):
            fill = C["accent_light"] if bt == "Classifier" else C["primary_light"]
            border = C["green_border"] if bt == "Classifier" else RGBColor(0xBF, 0xDB, 0xFE)
            self._rect(slide, x, y + Inches(0.35), bw, Inches(0.72), fill, border)
            self._textbox(slide, x, y + Inches(0.42), bw, Inches(0.22), bt, 10, True, C["primary_dark"] if bt != "Classifier" else C["accent"], PP_ALIGN.CENTER)
            self._textbox(slide, x, y + Inches(0.62), bw, Inches(0.22), bd, 9, False, C["text_sec"], PP_ALIGN.CENTER)
            if i < len(blocks) - 1:
                self._textbox(slide, x + bw, y + Inches(0.55), Inches(0.28), Inches(0.3), "→", 14, False, C["text_muted"], PP_ALIGN.CENTER)
            x += bw + Inches(0.28)
        self._h3(slide, ML, y + Inches(1.2), CW, "각 블록 내부 구성")
        col_w = (CW - Inches(0.25)) / 2
        self._code_block(slide, ML, y + Inches(1.52), col_w, Inches(1.35), [
            "# Conv Block 구성 흐름",
            "Conv2d(3×3, padding=1)",
            "→ BatchNorm2d → ReLU",
            "→ MaxPool2d(2×2)",
        ])
        self._code_block(slide, ML + col_w + Inches(0.25), y + Inches(1.52), col_w, Inches(1.35), [
            "# Classifier 구성 흐름",
            "Flatten",
            "→ Linear(50176, 512)",
            "→ ReLU → Dropout(0.5)",
            "→ Linear(512, 38)",
        ])
        self._insight(slide, ML, y + Inches(3.05), CW, Inches(0.75),
                      "채널 수를 32→64→128→256으로 점진적 확장하며 고수준 의미론적 특징을 추출하고, MaxPool로 공간 차원을 절반씩 축소하여 계산 효율성을 확보합니다.")

    def _slide_07(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "3-2. 아키텍처 최적화: 5가지 모델 비교 분석",
                         "자체 CNN → GAP 경량화 → 전이학습 모델(ResNet50, MobileNetV3, EfficientNet-B0) 비교")
        # timeline
        ty = y
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML + Inches(0.12), ty, Pt(3), Inches(1.55))
        line.fill.solid()
        line.fill.fore_color.rgb = C["border"]
        line.line.fill.background()
        steps = [
            ("Step 1. 전통적 CNN의 한계 발견 (Flatten 사용)",
             "Flatten 층으로 파라미터 약 2,610만 개 폭발. ResNet50(2,350만)보다 무거워 모바일 배포 부적합."),
            ("Step 2. 자체 경량화 달성 (GAP 도입) → 파라미터 98% 감소",
             "GAP(Global Average Pooling) 교체. 2,610만 → 54만 개(0.5M)로 98% 감소."),
            ("Step 3. 다양한 전이학습 모델 비교 실험",
             "ResNet50, MobileNetV3, EfficientNet-B0 3가지 SOTA 모델 성능·효율 비교."),
        ]
        for i, (st, sd) in enumerate(steps):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, ML + Inches(0.06), ty + Inches(0.05) + i * Inches(0.52), Inches(0.14), Inches(0.14))
            dot.fill.solid()
            dot.fill.fore_color.rgb = C["primary"]
            dot.line.color.rgb = C["primary_light"]
            self._textbox(slide, ML + Inches(0.35), ty + i * Inches(0.52), CW - Inches(0.35), Inches(0.22), st, 12, True)
            self._textbox(slide, ML + Inches(0.35), ty + Inches(0.22) + i * Inches(0.52), CW - Inches(0.35), Inches(0.28), sd, 10, False, C["text_sec"])
        self._table(slide, ML, y + Inches(1.75), CW, [
            ["모델", "분류기 구조", "파라미터 수", "평가"],
            ["Baseline CNN", "Flatten → Dense", "26,099,494", "가장 무거움"],
            ["ResNet50", "GAP + Residual Connection", "23,585,894", "깊은 층, 안정적 학습"],
            ["MobileNetV3", "GAP + Depthwise Conv", "4,250,710", "모바일 배포 최적화"],
            ["EfficientNet-B0", "GAP + Compound Scaling", "4,056,226", "최적 효율 + 99.73%"],
            ["Baseline CNN (GAP)", "GAP → Dense", "540,454", "극강 초경량 (98%↓)"],
        ], [Inches(2.2), Inches(3.8), Inches(2.3), Inches(3.2)], row_h=Inches(0.33))

    def _slide_08(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "4. 학습 설정", "하이퍼파라미터, 손실 함수 및 최적화 전략")
        self._table(slide, ML, y, CW, [
            ["항목", "설정값", "선택 이유"],
            ["프레임워크", "PyTorch 2.x + CUDA", "GPU 가속 병렬 연산"],
            ["손실 함수", "CrossEntropyLoss", "다중 분류(38 클래스) 최적 손실 함수"],
            ["옵티마이저", "Adam", "적응적 학습률, 안장점 돌파 능력"],
            ["학습률 (LR)", "0.001", "Adam 기본 권장값, 안정적 수렴"],
            ["배치 사이즈", "32", "GPU 메모리 대비 효율적 크기"],
            ["에포크", "10", "베이스라인 성능 확인용 초기 실험"],
            ["Dropout", "0.5", "완전 연결 계층의 과적합 방지"],
        ], [Inches(2.0), Inches(3.2), Inches(6.3)])
        self._h3(slide, ML, y + Inches(2.95), CW, "가중치 저장 전략")
        self._card(slide, ML, y + Inches(3.28), Inches(5.95), Inches(0.95), "best_model.pth",
                   "검증 정확도(Val Acc) 최고치를 갱신할 때마다 자동 저장 → 과적합 발생 시에도 최적 시점의 가중치 확보")
        self._card(slide, ML + Inches(6.15), y + Inches(3.28), Inches(5.95), Inches(0.95), "last_model.pth",
                   "마지막 에포크 종료 후 저장 → 학습 재개(Resume) 및 추가 학습(Fine-tuning)에 활용")

    def _slide_09(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "5. 실험실 모델 성과 (EfficientNet-B0)", "PlantVillage 데이터셋 기반 10 Epoch 학습 결과")
        bw = (CW - Inches(0.36)) / 3
        self._result_box(slide, ML, y, bw, Inches(1.05), "Best Val Loss", "0.0094", C["primary"])
        self._result_box(slide, ML + bw + Inches(0.18), y, bw, Inches(1.05), "Best Val Accuracy", "99.73%", C["accent"], C["green_bg"], C["green_border"])
        self._result_box(slide, ML + 2 * (bw + Inches(0.18)), y, bw, Inches(1.05), "Best Epoch", "9", C["warning"])
        self._insight(slide, ML, y + Inches(1.2), CW, Inches(0.85),
                      "완벽한 실험실 AI 모델 구축 성공: EfficientNet-B0은 효율적인 파라미터(405만 개)로 38개 클래스를 거의 완벽하게 분류합니다.",
                      "success")
        self._h3(slide, ML, y + Inches(2.15), CW, "그런데... 이 모델이 실제 야외 환경에서도 통할까?")
        self._insight(slide, ML, y + Inches(2.5), CW, Inches(0.7),
                      "실험실의 단색 배경에서 99.73%를 달성한 이 모델을, 흙과 잡초가 뒤섞인 실제 야외 환경(PlantDoc 데이터)에 투입해 보았습니다...",
                      "danger")

    def _slide_10(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "6. 치명적 문제 발생: 도메인 시프트", "실험실에서 야외(In-the-wild)로 나갔을 때 벌어진 충격적인 결과")
        bw = (CW - Inches(0.72)) / 3
        self._result_box(slide, ML, y, bw, Inches(1.1), "실험실 정확도", "99.73%", C["accent"], C["green_bg"], C["green_border"])
        self._textbox(slide, ML + bw + Inches(0.18), y + Inches(0.35), Inches(0.5), Inches(0.5), "→", 36, True, C["text"], PP_ALIGN.CENTER)
        self._result_box(slide, ML + bw + Inches(0.72), y, bw, Inches(1.1), "야외 테스트 정확도", "19.12%", C["danger"], C["red_bg"], C["red_border"])
        self._h3(slide, ML, y + Inches(1.25), CW, "무엇이 문제인가?")
        self._bullets(slide, ML, y + Inches(1.55), CW, Inches(1.35), [
            "실험실 데이터(PlantVillage): 단색 배경에 잎사귀 하나만 깨끗하게 촬영된 이미지",
            "야외 데이터(PlantDoc): 흙, 잡초, 가지, 그림자 등 복잡한 배경이 뒤섞인 실제 환경 이미지",
            "단색 배경 패턴에 완전히 과적합 → 전형적인 도메인 시프트(Domain Shift) 현상 발생",
        ])
        self._insight(slide, ML, y + Inches(3.05), CW, Inches(0.75),
                      "핵심 발견: 실험실에서 아무리 높은 정확도를 달성해도, 실제 현장 환경과의 괴리(Domain Gap)를 해결하지 못하면 실무 적용이 불가능하다는 사실을 확인했습니다.",
                      "danger")

    def _slide_11(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "7. 원인 분석: Grad-CAM 시각화", "AI가 대체 어디를 보고 판단하는지, 모델의 시선을 직접 시각화")
        lw, rw = Inches(5.9), Inches(5.9)
        self._h3(slide, ML, y, lw, "Grad-CAM이란?")
        self._textbox(slide, ML, y + Inches(0.32), lw, Inches(0.55),
                      "Gradient-weighted Class Activation Mapping — 예측 시 이미지의 어느 부위에 집중했는지 히트맵으로 시각화하는 XAI 기법입니다.",
                      11, False, C["text_sec"])
        self._h3(slide, ML, y + Inches(0.95), lw, "분석 결과")
        self._bullets(slide, ML, y + Inches(1.25), lw, Inches(1.5), [
            "발견: AI는 잎사귀의 병변을 전혀 보지 않음",
            "오류 원인: 뒷배경의 흙과 나뭇가지에 시선이 집중(활성화)",
            "예시: Apple Scab → Peach Bacterial Spot 오진",
        ], 11)
        self._image_card(slide, ML + lw + Inches(0.35), y, rw, Inches(2.85),
                         "gradcam_results/wrong_2_true_Apple___Apple_scab_pred_Peach___Bacterial_spot.png",
                         "Grad-CAM 히트맵: 잎이 아닌 흙과 배경에 시선(붉은색)이 쏠림")
        self._insight(slide, ML, y + Inches(3.0), CW, Inches(0.75),
                      "근본 원인 확정: 모델이 잎의 병변이 아닌 흙이나 배경을 보고 판단하는 치명적 과적합. 3가지 가설 검증 실험에 돌입.",
                      "danger")

    def _slide_12(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "8-1. 가설 1: 물리적 전처리 도입  [실패]",
                         "배경이 문제라면, AI(rembg)로 배경을 물리적으로 제거하면 해결될까?")
        lw = Inches(5.9)
        self._h3(slide, ML, y, lw, "실험 설계")
        self._bullets(slide, ML, y + Inches(0.32), lw, Inches(0.7), [
            "방법: rembg로 야외 데이터 전체 배경 제거(누끼)",
            "기대: 배경 노이즈 제거 → 잎사귀 집중",
        ], 11)
        self._result_box(slide, ML, y + Inches(1.05), Inches(3.5), Inches(0.95), "결과 정확도", "47.81%", C["danger"], C["danger_light"], C["red_border"])
        self._h3(slide, ML, y + Inches(2.1), lw, "실패 원인 분석")
        self._bullets(slide, ML, y + Inches(2.4), lw, Inches(0.75), [
            "환경적 문맥(Context)까지 파괴됨",
            "잎 끝부분(병변)까지 배경으로 인식되어 훼손",
        ], 11)
        self._image_card(slide, ML + lw + Inches(0.35), y, Inches(5.9), Inches(2.9),
                         "assets/curve_v3.png", "가설 1 (배경 제거) 학습 곡선: 47.8%에서 수렴")
        self._insight(slide, ML, y + Inches(3.25), CW, Inches(0.6),
                      "결론: 물리적 전처리(배경 제거)는 환경 문맥을 상실시키고 병변을 훼손하여, 오히려 모델의 판단 근거를 약화시킵니다.",
                      "danger")

    def _slide_13(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "8-2. 가설 2: 혼합 도메인 학습  [실패]",
                         "기존 실험실 데이터와 야외 데이터를 섞어서 학습하면 해결될까?")
        lw = Inches(5.9)
        self._h3(slide, ML, y, lw, "실험 설계")
        self._bullets(slide, ML, y + Inches(0.32), lw, Inches(0.7), [
            "방법: 실험실 5,000장 + 야외 2,500장 혼합 파인튜닝",
            "기대: 깨끗한 데이터로 질병 패턴 망각 방지",
        ], 11)
        self._result_box(slide, ML, y + Inches(1.05), Inches(3.5), Inches(0.95), "결과 정확도", "47.01%", C["danger"], C["danger_light"], C["red_border"])
        self._h3(slide, ML, y + Inches(2.1), lw, "실패 원인 분석")
        self._bullets(slide, ML, y + Inches(2.4), lw, Inches(0.75), [
            "쉬운 데이터가 섞이면서 모델의 시선이 분산",
            "야외 환경 적응에 대한 학습 집중도 희석",
        ], 11)
        self._image_card(slide, ML + lw + Inches(0.35), y, Inches(5.9), Inches(2.9),
                         "assets/curve_v4.png", "가설 2 (데이터 혼합) 학습 곡선: 47.0%에서 한계")
        self._insight(slide, ML, y + Inches(3.25), CW, Inches(0.6),
                      "결론: 쉬운 데이터를 혼합하면 야외 환경 적응에 대한 모델의 시선이 분산되어, 도메인 적응에 오히려 역효과를 초래합니다.",
                      "danger")

    def _slide_14(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "8-3. 가설 3: 본질적 도메인 적응  [성공]",
                         "원본 야외 데이터만 투입 + 가혹한 증강 + 최신 수학적 최적화로 전이학습")
        self._h3(slide, ML, y, CW, "실험 설계")
        cw = (CW - Inches(0.36)) / 3
        cards = [
            ("타겟 도메인 집중", "야외 데이터(PlantDoc) 원본만 집중 투입하여 전이학습 수행", C["primary"]),
            ("가혹한 데이터 증강", "ColorJitter, RandomResizedCrop 등 극한 증강 적용", C["warning"]),
            ("수학적 최적화", "Cosine Annealing LR + Label Smoothing 기법 도입", C["accent"]),
        ]
        for i, (t, d, ac) in enumerate(cards):
            self._card(slide, ML + i * (cw + Inches(0.18)), y + Inches(0.32), cw, Inches(0.95), t, d, accent_left=ac)
        ly, lx = y + Inches(1.45), ML
        lw, rw = Inches(5.9), Inches(5.9)
        bw = (lw - Inches(0.36)) / 3
        self._result_box(slide, lx, ly, bw, Inches(0.95), "기존 야외 정확도", "19.1%", C["danger"])
        self._textbox(slide, lx + bw + Inches(0.05), ly + Inches(0.25), Inches(0.4), Inches(0.4), "→", 28, True, C["text"], PP_ALIGN.CENTER)
        self._result_box(slide, lx + bw + Inches(0.5), ly, bw, Inches(0.95), "최종 정확도", "60.6%", C["accent"], C["green_bg"], C["green_border"])
        self._insight(slide, lx, ly + Inches(1.05), lw, Inches(0.7),
                      "3배 이상 향상! (배경 제거, 데이터 혼합)를 버리고 강건성 최적화에 집중한 방법이 가장 효과적임을 입증했습니다.",
                      "success")
        self._image_card(slide, lx + lw + Inches(0.35), ly, rw, Inches(1.85),
                         "assets/curve_v2.png", "가설 3 학습 곡선: 60.56%까지 지속 상승")

    def _slide_15(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "9. 결론: 가설 검증 종합 비교", "3가지 가설의 실험 결과를 정량적으로 비교 분석")
        self._table(slide, ML, y, CW, [
            ["가설", "전략", "정확도", "결과"],
            ["기존 (도메인 시프트)", "실험실 가중치 그대로 야외 테스트", "19.12%", "실패"],
            ["가설 1 (물리적 전처리)", "rembg로 배경 제거 후 파인튜닝", "47.81%", "실패"],
            ["가설 2 (혼합 도메인)", "실험실+야외 데이터 섞어서 학습", "47.01%", "실패"],
            ["가설 3 (본질적 적응)", "가혹한 증강 + 수학적 최적화 전이학습", "60.56%", "최선"],
        ], [Inches(2.5), Inches(4.8), Inches(1.8), Inches(2.5)])
        self._insight(slide, ML, y + Inches(2.05), CW, Inches(1.1),
                      "핵심 교훈 및 성과: 방대한 실험실 데이터로 학습된 기존 지식(가중치)을 적극 활용한 전이학습 덕분에, 단 2,589장의 매우 적은 야외 데이터만으로도 60.56%라는 유의미한 도메인 적응 성과를 달성했습니다. 소량의 타겟 데이터 + 강력한 증강 기법 결합 파인튜닝이 실무적으로 효율적인 해결책임을 증명합니다.",
                      "success")

    def _slide_16(self):
        slide = self._blank()
        self._progress(slide)
        y = self._header(slide, "10. 향후 연구 방향", "도메인 적응 성능(80~90% 이상)을 목표로 한 고도화 전략")
        cards = [
            ("Vision Transformer (ViT) 도입", "Self-Attention으로 배경 노이즈 무시, 잎사귀 병변 집중", C["primary"]),
            ("적대적 도메인 적응 (DANN)", "Gradient Reversal Layer로 도메인 차이 제거", C["accent"]),
            ("객체 탐지(YOLO) 파이프라인 결합", "잎사귀 Bounding Box 크롭 후 분류로 노이즈 차단", C["warning"]),
            ("자기 지도 학습 (Self-Supervised)", "야외 식물 사진 Pre-training 후 파인튜닝", C["danger"]),
        ]
        gw, gh, gap = Inches(5.95), Inches(1.25), Inches(0.18)
        for i, (t, d, ac) in enumerate(cards):
            col, row = i % 2, i // 2
            self._card(slide, ML + col * (gw + gap), y + row * (gh + gap), gw, gh, t, d, accent_left=ac)
        self._insight(slide, ML, y + Inches(2.75), CW, Inches(0.75),
                      "방향성: 현재 V2 모델(60.56%)의 아키텍처를 기반으로, 위 전략들을 단계적으로 적용하여 야외 환경에서의 정확도를 80~90% 이상으로 끌어올리는 것을 목표로 후속 연구를 진행할 계획입니다.")

    def _slide_17(self):
        self._cover(
            "감사합니다",
            "질의응답 (Q&A)",
            "딥러닝 기반 농작물 잎 이미지 병해충 진단 시스템\n도메인 시프트 극복 연구 보고",
            variant="end",
        )


if __name__ == "__main__":
    out = Builder().build()
    print(f"Created: {out}")
    print("Slides: 17")
